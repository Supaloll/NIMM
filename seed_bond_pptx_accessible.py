"""
seed_bond_pptx_accessible.py — Bond "Créer une présentation PowerPoint accessible"

Prérequis :
    pip install python-pptx --break-system-packages
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from core.database import save_prompt, list_prompts, get_all_users, set_user_context
except ImportError as e:
    print(f"Erreur d'import : {e}")
    sys.exit(1)

_users = get_all_users()
if not _users:
    print("Aucun utilisateur trouvé. Lancez NIMM au moins une fois.")
    sys.exit(1)
set_user_context(_users[0]["id"])
print(f"Contexte utilisateur : {_users[0].get('name', _users[0]['id'])}")

LABEL = "Créer une présentation PowerPoint accessible (PPTX)"

TEXT = """Quand utiliser ce bond : quand l'utilisateur veut créer une présentation
PowerPoint qui respecte les critères d'accessibilité — lisible par les lecteurs
d'écran, avec alt text sur les images, titres de diapositive obligatoires,
ordre de lecture logique, contraste suffisant et langue déclarée.

Prérequis (une seule fois) :
    pip install python-pptx --break-system-packages

Méthode :

1. Recueillir le contenu de la présentation :
   - Titre général et auteur
   - Pour chaque diapositive : titre, corps du texte, images éventuelles avec
     leur description (alt text), notes du présentateur si souhaitées
   - Langue principale du document (fr par défaut)

2. Créer la présentation avec python-pptx :
   from pptx import Presentation
   from pptx.util import Inches, Pt
   from pptx.dml.color import RGBColor
   prs = Presentation()

3. Règles d'accessibilité OBLIGATOIRES pour chaque diapositive :

   a) Titre présent et non vide :
      slide = prs.slides.add_slide(prs.slide_layouts[1])  # layout titre+contenu
      slide.shapes.title.text = "Titre de la diapositive"
      — Ne jamais laisser un titre vide ou remplacé par une image.

   b) Ordre de lecture logique :
      Les formes (shapes) sont lues dans l'ordre du panneau de sélection.
      Ajouter les formes dans l'ordre logique : titre d'abord, puis corps.

   c) Alt text sur toutes les images :
      from pptx.oxml.ns import qn
      pic = slide.shapes.add_picture(chemin_image, left, top, width, height)
      pic._element.nvPicPr.cNvPr.set('descr', "Description textuelle de l'image")
      pic._element.nvPicPr.cNvPr.set('title', "Titre court de l'image")

   d) Contraste de texte suffisant (ratio 4.5:1 minimum) :
      Utiliser noir (#000000) sur blanc (#FFFFFF) ou équivalent fort.
      Éviter gris clair sur blanc, jaune sur blanc, etc.

   e) Ne pas transmettre d'information par la couleur seule :
      Si un tableau a des cellules colorées pour signifier "OK/KO",
      ajouter aussi le texte "OK" ou "KO" dans la cellule.

   f) Langue du document déclarée dans le XML core properties :
      from pptx.opc.constants import RELATIONSHIP_TYPE as RT
      prs.core_properties.language = "fr-FR"

   g) Notes du présentateur pour chaque diapositive :
      slide.notes_slide.notes_text_frame.text = "Description audio de la diapositive"
      Les notes servent de description pour les participants non-voyants.

4. Police recommandée : Arial, Calibri ou Verdana — éviter les polices
   manuscrites ou décoratives. Taille minimale : 18pt pour le corps, 24pt pour
   les titres.

5. Sauvegarder :
   prs.save(chemin_sortie)
   print(f"Présentation enregistrée : {chemin_sortie}")

6. Après la création, lancer une vérification rapide :
   — Toutes les diapositives ont-elles un titre ?
   — Toutes les images ont-elles un alt text non vide ?
   — Y a-t-il des notes du présentateur sur chaque diapositive ?
   Afficher un rapport "Vérification accessibilité : X/Y diapositives conformes."

Règles importantes :
- Ne jamais créer de diapositive "titre seul" avec tout le texte dans une zone
  flottante non liée au layout — cela brise l'ordre de lecture des lecteurs d'écran.
- Les animations et transitions ne sont pas accessibles : les éviter ou les réduire
  au minimum. Les signaler si l'utilisateur en demande.
- Pour les tableaux dans les diapositives : toujours définir la première ligne
  comme ligne d'en-tête via table.rows[0].cells[0]._tc.get_or_add_tcPr().
"""

META = {
    "description": "Créer une présentation PowerPoint (.pptx) conforme aux critères d'accessibilité : titres, alt text, ordre de lecture, contraste.",
    "mots_cles": [
        "powerpoint", "pptx", "présentation", "diaporama", "diapositive",
        "accessible", "accessibilité", "alt text", "lecteur écran", "contraste",
        "titre", "python-pptx", "slide", "créer", "générer",
    ],
    "valide": True,
    "valide_par_laurent": True,
    "version": 1,
}


def main():
    existing = list_prompts("skill")
    for sid, sk in (existing.items() if isinstance(existing, dict) else {}):
        lbl = (sk.get("label") or "").lower()
        if "powerpoint" in lbl or ("pptx" in lbl and "accessible" in lbl):
            print(f"Bond similaire déjà présent (id={sid}) : {sk.get('label')}")
            return
    entry = save_prompt(None, LABEL, TEXT, type="skill", meta=META)
    if entry:
        print(f"Bond créé : {LABEL}")
    else:
        print("Erreur lors de la création.")


if __name__ == "__main__":
    main()
