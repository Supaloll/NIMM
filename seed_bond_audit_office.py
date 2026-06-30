"""
seed_bond_audit_office.py — Bond "Audit accessibilité document Office (PPTX/DOCX)"

Prérequis :
    pip install python-pptx python-docx --break-system-packages
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from core.database import save_prompt, list_prompts, get_all_users, set_user_context
except ImportError as e:
    print(f"Erreur d'import : {e}")
    sys.exit(1)

_users = get_all_users()
if not _users:
    print("Aucun utilisateur trouvé. Lancez NIMM au moins une fois.")
    sys.exit(1)
set_user_context(_users[0]["id"])
print(f"Contexte utilisateur : {_users[0].get('name', _users[0]['id'])}")

LABEL = "Audit d'accessibilité document Office (PPTX et DOCX)"

TEXT = """Quand utiliser ce bond : quand l'utilisateur veut vérifier si un fichier
PowerPoint (.pptx) ou Word (.docx) existant respecte les critères d'accessibilité
avant de le diffuser. Produit un rapport clair avec les problèmes trouvés et
les corrections à apporter.

Prérequis (une seule fois) :
    pip install python-pptx python-docx --break-system-packages

Méthode — Audit PPTX :

1. Ouvrir le fichier :
   from pptx import Presentation
   from pptx.oxml.ns import qn
   prs = Presentation(chemin_pptx)

2. Vérifier chaque diapositive :
   problemes = []
   for i, slide in enumerate(prs.slides, 1):
       # a) Titre manquant
       titre = slide.shapes.title
       if not titre or not (titre.text or '').strip():
           problemes.append(f"Diapo {i} : TITRE MANQUANT ou vide")

       # b) Alt text manquant sur les images
       for shape in slide.shapes:
           if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
               descr = shape._element.nvPicPr.cNvPr.get('descr', '')
               if not descr or not descr.strip():
                   problemes.append(f"Diapo {i} : image '{shape.name}' sans alt text")

       # c) Notes du présentateur absentes
       notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ''
       if not notes.strip():
           problemes.append(f"Diapo {i} : notes du présentateur vides (description audio manquante)")

       # d) Texte trop petit (< 18pt)
       for shape in slide.shapes:
           if shape.has_text_frame:
               for para in shape.text_frame.paragraphs:
                   for run in para.runs:
                       taille = run.font.size
                       if taille and taille.pt < 18:
                           problemes.append(f"Diapo {i} : texte trop petit ({taille.pt:.0f}pt) dans '{shape.name}'")

3. Afficher le rapport :
   if not problemes:
       print("Aucun problème d'accessibilité détecté dans la présentation.")
   else:
       print(f"Rapport d'accessibilité PPTX : {len(problemes)} problème(s) trouvé(s)")
       for p in problemes:
           print(f"  - {p}")

Méthode — Audit DOCX :

1. Ouvrir le fichier :
   from docx import Document
   from docx.oxml.ns import qn
   doc = Document(chemin_docx)

2. Vérifications :
   problemes = []

   # a) Langue du document
   lang = doc.core_properties.language
   if not lang:
       problemes.append("Langue du document non déclarée dans les propriétés")

   # b) Hiérarchie des titres (pas de saut de niveau h1 → h3)
   niveaux = []
   for para in doc.paragraphs:
       if para.style.name.startswith('Heading'):
           try:
               n = int(para.style.name.split(' ')[-1])
               if niveaux and n > niveaux[-1] + 1:
                   problemes.append(f"Saut de niveau de titre : de Titre {niveaux[-1]} à Titre {n} ('{para.text[:40]}')")
               niveaux.append(n)
           except ValueError:
               pass

   # c) Images sans alt text
   for rel in doc.part.rels.values():
       if 'image' in rel.reltype:
           # Chercher la description dans le XML de l'image
           pass  # python-docx ne donne pas accès direct à l'alt text facilement
   # Signaler la limitation : vérifier les alt text manuellement dans Word
   problemes.append("Note : vérifier manuellement les alt text des images dans Word (Format > Alt Text)")

   # d) Tableaux sans en-têtes
   for i, table in enumerate(doc.tables, 1):
       premiere_ligne_est_entete = False
       for cell in table.rows[0].cells:
           if cell.paragraphs[0].style.name in ('Table Header', 'En-tête de tableau'):
               premiere_ligne_est_entete = True
               break
       if not premiere_ligne_est_entete:
           problemes.append(f"Tableau {i} : première ligne non marquée comme en-tête")

3. Afficher le rapport DOCX avec les mêmes conventions.

Format du rapport final :
   "=== Rapport d'accessibilité : [nom du fichier] ==="
   "Type : PPTX / DOCX"
   "Problèmes critiques (bloquants pour lecteur d'écran) : X"
   "Problèmes mineurs (améliorations recommandées) : Y"
   [Liste des problèmes]
   "=== Fin du rapport ==="

Règles importantes :
- Classer les problèmes par sévérité : titre manquant et alt text manquant
  sont critiques ; taille de texte et notes sont des recommandations.
- Proposer après le rapport de corriger automatiquement les problèmes
  récupérables (ex. ajouter des notes vides, normaliser les tailles de police).
- Ne jamais modifier le fichier original sans demande explicite — toujours
  proposer de sauvegarder dans un nouveau fichier avec suffixe "_accessible".
"""

META = {
    "description": "Auditer l'accessibilité d'un fichier PPTX ou DOCX existant et produire un rapport des problèmes à corriger.",
    "mots_cles": [
        "audit", "accessibilité", "pptx", "docx", "word", "powerpoint",
        "vérifier", "contrôle", "rapport", "problèmes", "titres",
        "alt text", "images", "lecteur écran", "conformité", "wcag", "rgaa",
    ],
    "valide": True,
    "valide_par_laurent": True,
    "version": 1,
}


def main():
    existing = list_prompts("skill")
    for sid, sk in (existing.items() if isinstance(existing, dict) else {}):
        lbl = (sk.get("label") or "").lower()
        if "audit" in lbl and ("office" in lbl or "pptx" in lbl or "docx" in lbl):
            print(f"Bond similaire déjà présent (id={sid}) : {sk.get('label')}")
            return
    entry = save_prompt(None, LABEL, TEXT, type="skill", meta=META)
    if entry:
        print(f"Bond créé : {LABEL}")
    else:
        print("Erreur lors de la création.")


if __name__ == "__main__":
    main()
